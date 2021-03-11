from pptx.text.text import _Run
from module_pptx.pptx_styles.PptxTextFrameStyle import PptxTextFrameStyle
from module_pptx.pptx_styles.PptxParagraphStyle import PptxParagraphStyle
from module_pptx.PptxTextContentEntry import PptxTextContentEntry


class PptxTextContent:
    def __init__(self, parameters, styles_input, text_frame):
        self._parameters = parameters
        self._styles_input: dict = None
        if styles_input is not None:
            self._styles_input = dict(styles_input)
        self._text_frame = text_frame
        self._paragraph = None
        self._text_frame_style: PptxTextFrameStyle = None
        self._paragraph_style: PptxParagraphStyle = None

    def update_styles_input(self, styles_input):
        if styles_input is not None:
            if self._styles_input is not None:
                self._styles_input.update(styles_input)
            else:
                self._styles_input = styles_input

    def apply_text_frame_styles(self):
        self._text_frame.clear()
        self._paragraph = self._text_frame.paragraphs[0]
        if self._styles_input is not None:
            self._text_frame_style = PptxTextFrameStyle(self._styles_input)
            self._text_frame_style.generate()
            self._text_frame_style.apply_styles(self._text_frame)

    def apply_paragraph_styles(self):
        if self._styles_input is not None:
            paragraph_style = PptxParagraphStyle(self._styles_input)
            paragraph_style.generate()
            paragraph_style.apply_styles(self._paragraph)

    def apply_content_entry_styles(self):
        for entry in self._parameters:
            run: _Run = self._paragraph.add_run()
            text_entry = PptxTextContentEntry(entry, self._styles_input, run)
            text_entry.generate()

    def apply_styles(self):
        self.apply_text_frame_styles()
        self.apply_paragraph_styles()
        self.apply_content_entry_styles()
