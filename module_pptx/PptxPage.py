from pptx import Presentation
from module_misc.BasicTask import BasicTask


class PptxPage(BasicTask):
    def __init__(self, parameters, presentation: Presentation, layout, settings=None):
        super().__init__(parameters, settings)
        self._presentation = presentation
        self._layout = layout
        self._slide = None
        self._instructions_input = None
        self._instructions = []

    def create_page(self):
        layout = self._presentation.slide_layouts[self._layout]
        self._slide = self._presentation.slides.add_slide(layout)
        for shape in self._slide.placeholders:
            print('%d %s' % (shape.placeholder_format.idx, shape.name))

    def generate(self):
        super().generate()
        self.create_page()
