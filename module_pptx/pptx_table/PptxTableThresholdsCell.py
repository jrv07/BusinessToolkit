from pptx.table import _Cell
from module_pptx.pptx_styles.PptxCellStyle import PptxCellStyle
from module_pptx.PptxTextContent import PptxTextContent
from module_pptx.pptx_table.PptxTableCellBase import PptxTableCellBase
import logging


class PptxTableThresholdsCell(PptxTableCellBase):
    def __init__(self, parameters, default_cell_style=None):
        super().__init__(parameters, default_cell_style)
        self._thresholds = None
        self._operator = None
        self._value = None
        self._styles = None
        self._styles_input = None
        self._value_format: str = None
        self._extractor.add_float("thresholds_float", specifier="thresholds", optional=False,
                                  alternative="thresholds_list",
                                  description="Can be used if there is exactly one threshold. Specify two styles.")
        self._extractor.add_list("thresholds_list", specifier="thresholds", entry_type=float, optional=False,
                                 description="Use it when there are multiple thresholds. "
                                             "the thresholds must be sorted by the used operator. For example: "
                                             "operator: '<', then the thresholds list is like [2, 3.5, 8]"
                                             "When there are n thresholds, specify n+1 styles.")
        self._extractor.add_float("value", description="The value displayed in the cell. Optional, because it "
                                                       "can be set later.")
        self._extractor.add_list("styles", entry_type=dict, entry_object_type=PptxCellStyle, optional=False,
                                 description="When there are n thresholds, specify n+1 styles. The styles list must be "
                                             "sorted in the same way as the thresholds. "
                                             "For example: thresholds list: [5, 10] => Specify styles [S0, S1, S2]. "
                                             "With operator: '<' the cell with value 11 would be of style S2.")
        valid_operators = ["gt", "lt", "gte", "lte", ">", "<", ">=", "<="]
        self._extractor.add_str("operator", optional=False, valid_values=valid_operators,
                                description="The operator that will be used in the evaluation of the value.")
        # Todo: Create a value format class with user friendly instructions
        self._extractor.add_str("value_format", description="This format is applied to the value before printing it to "
                                                            "the cell. For example '{:.2f}' rounds the value to 2 "
                                                            "decimals.")
        self._extractor.set_group_name("Presentation")

    def set_cell_value(self, value: float):
        self._value = value

    def generate_styles(self):
        styles = []
        for style_input in self._styles_input:
            style = PptxCellStyle(self._default_cell_style)
            style.update_styles_by_input(style_input)
            style.generate()
            styles.append(style)
        self._styles = styles

    def apply_format(self, value):
        if self._value_format is not None:
            return self._value_format.format(value)
        else:
            return str(value)

    def apply_text_frame_style(self, style_input, text_frame):
        content = [
            {
                "text": self.apply_format(self._value),
                "font": style_input
            }
        ]
        text_content = PptxTextContent(content, self._default_cell_style, text_frame)
        text_content.apply_styles()

    @staticmethod
    def evaluate_gt(value1, value2) -> bool:
        if value1 > value2:
            return True
        else:
            return False

    @staticmethod
    def evaluate_lt(value1, value2) -> bool:
        if value1 < value2:
            return True
        else:
            return False

    @staticmethod
    def evaluate_gte(value1, value2) -> bool:
        if value1 >= value2:
            return True
        else:
            return False

    @staticmethod
    def evaluate_lte(value1, value2) -> bool:
        if value1 <= value2:
            return True
        else:
            return False

    def evaluate(self, value1, value2) -> bool:
        if self._operator in ["gt", ">"]:
            return self.evaluate_gt(value1, value2)
        elif self._operator in ["lt", "<"]:
            return self.evaluate_lt(value1, value2)
        elif self._operator in ["gte", ">="]:
            return self.evaluate_gte(value1, value2)
        elif self._operator in ["lte", "<="]:
            return self.evaluate_lte(value1, value2)

    def evaluate_value(self):
        for index, threshold in enumerate(self._thresholds):
            if self.evaluate(self._value, threshold):
                return index
        return len(self._thresholds)

    def apply_styles(self, cell: _Cell):
        style_index = self.evaluate_value()
        self._styles[style_index].apply_styles(cell)
        self.apply_text_frame_style(self._styles_input[style_index], cell.text_frame)

    def get_thresholds(self):
        threshold = self._extractor.get_value("thresholds_float")
        if threshold is None:
            self._thresholds = self._extractor.get_value("thresholds_list")
            last_threshold = None
            for threshold in self._thresholds:
                if last_threshold is not None:
                    if not self.evaluate(last_threshold, threshold):
                        logging.error("thresholds list must be given sorted by operator. Current values: {}"
                                      .format(self._thresholds))
                        raise ValueError
                last_threshold = threshold
        else:
            self._thresholds = [threshold]

    def get_styles(self):
        self._styles_input = self._extractor.get_value("styles")
        if len(self._styles_input) != len(self._thresholds)+1:
            logging.error("You must give the correct number of Cell Styles. {} styles are needed, current count: {}"
                          .format(len(self._thresholds)+1, len(self._styles_input)))
            raise ValueError

    def get_arguments(self):
        super().get_arguments()
        self._operator = self._extractor.get_value("operator")
        self.get_thresholds()
        self._value = self._extractor.get_value("value")
        self.get_styles()
        self._value_format = self._extractor.get_value("value_format")

    def generate(self):
        self.extract()
        self.get_arguments()
        self.generate_styles()
